"""Deterministic fixture generation.

Fixtures are *built*, never committed: that keeps the repository free of any
third-party or sensitive document, and guarantees every file is small and legally
unambiguous. Generation is idempotent -- rerunning overwrites with identical bytes
for everything except the OOXML formats, whose zip members carry timestamps.

Run standalone with ``python tests/fixtures/generate.py`` to inspect the output.
"""

from __future__ import annotations

from pathlib import Path

import fitz  # PyMuPDF
from PIL import Image, ImageDraw, ImageFont

GENERATED = Path(__file__).parent / "generated"

TERMINAL_TEXT = [
    "$ npm run build",
    "",
    "> media-dashboard@1.4.2 build",
    "> tsc -p tsconfig.json",
    "",
    "src/services/report.ts(84,17): error TS2345: Argument of type",
    "  'string | undefined' is not assignable to parameter of type 'string'.",
    "  Type 'undefined' is not assignable to type 'string'.",
    "",
    "Found 1 error in src/services/report.ts:84",
    "",
    "npm ERR! code ELIFECYCLE",
    "npm ERR! errno 2",
]

PDF_TEXT_PAGES = [
    (
        "Deployment Requirements",
        [
            "1. Runtime: Node.js 20 LTS or newer.",
            "2. The service must run behind TLS termination.",
            "3. Database migrations run before the new version starts.",
            "4. Health check endpoint: GET /healthz, expected 200 within 5s.",
        ],
    ),
    (
        "Rollback Procedure",
        [
            "If the health check fails twice, the deploy is rolled back.",
            "Rollback restores the previous container image tag.",
            "Database migrations are forward-only and are never reverted.",
        ],
    ),
]


def _font(size: int = 14) -> ImageFont.ImageFont:
    """A monospace font if the platform has one, else Pillow's bitmap default."""
    for candidate in ("DejaVuSansMono.ttf", "consola.ttf", "cour.ttf", "DejaVuSans.ttf"):
        try:
            return ImageFont.truetype(candidate, size)
        except OSError:
            continue
    return ImageFont.load_default()


def make_terminal_screenshot(path: Path) -> Path:
    """A dark terminal showing a TypeScript build failure."""
    width, height = 860, 300
    image = Image.new("RGB", (width, height), (18, 18, 22))
    draw = ImageDraw.Draw(image)
    font = _font(15)
    y = 14
    for line in TERMINAL_TEXT:
        colour = (235, 235, 235)
        if "error" in line or "ERR!" in line:
            colour = (255, 110, 110)
        elif line.startswith("$"):
            colour = (120, 230, 140)
        draw.text((14, y), line, fill=colour, font=font)
        y += 20
    image.save(path, format="PNG")
    return path


def make_ui_screenshot(path: Path) -> Path:
    """A minimal settings dialog with a disabled button and an inline error."""
    width, height = 640, 400
    image = Image.new("RGB", (width, height), (245, 246, 248))
    draw = ImageDraw.Draw(image)
    title_font = _font(20)
    font = _font(14)

    draw.rectangle([40, 40, 600, 360], fill=(255, 255, 255), outline=(206, 210, 216))
    draw.rectangle([40, 40, 600, 88], fill=(240, 242, 245), outline=(206, 210, 216))
    draw.text((60, 56), "Notification Settings", fill=(24, 26, 30), font=title_font)

    draw.text((60, 116), "Email address", fill=(60, 64, 70), font=font)
    draw.rectangle([60, 138, 580, 174], fill=(255, 255, 255), outline=(220, 60, 60), width=2)
    draw.text((70, 148), "not-an-email", fill=(40, 42, 48), font=font)
    draw.text((60, 182), "Enter a valid email address.", fill=(200, 40, 40), font=font)

    draw.rectangle([60, 220, 78, 238], fill=(255, 255, 255), outline=(120, 124, 130))
    draw.text((90, 220), "Send weekly digest", fill=(60, 64, 70), font=font)
    draw.rectangle([60, 252, 78, 270], fill=(45, 110, 220), outline=(45, 110, 220))
    draw.text((90, 252), "Notify on build failure", fill=(60, 64, 70), font=font)

    draw.rectangle([400, 300, 480, 336], fill=(255, 255, 255), outline=(180, 184, 190))
    draw.text((424, 310), "Cancel", fill=(90, 94, 100), font=font)
    draw.rectangle([496, 300, 580, 336], fill=(200, 205, 212), outline=(200, 205, 212))
    draw.text((524, 310), "Save", fill=(140, 144, 150), font=font)

    image.save(path, format="PNG")
    return path


def make_text_pdf(path: Path) -> Path:
    """A two-page PDF with a real, selectable text layer."""
    document = fitz.open()
    for heading, lines in PDF_TEXT_PAGES:
        page = document.new_page(width=595, height=842)
        page.insert_text((72, 92), heading, fontsize=18, fontname="helv")
        y = 132
        for line in lines:
            page.insert_text((72, y), line, fontsize=11, fontname="helv")
            y += 22
    document.save(path)
    document.close()
    return path


def make_scanned_pdf(path: Path) -> Path:
    """A one-page PDF whose only content is a raster image -- no text layer at all."""
    width, height = 992, 1403  # A4 at ~120 dpi; large enough for OCR, small on disk
    image = Image.new("RGB", (width, height), (252, 252, 250))
    draw = ImageDraw.Draw(image)
    draw.text((120, 160), "INVOICE 2026-0417", fill=(20, 20, 20), font=_font(46))
    body = [
        "Bill to: Northwind Analytics Ltd",
        "Issue date: 2026-03-11",
        "Due date: 2026-04-10",
        "",
        "Line item              Qty      Amount",
        "Support retainer         1      1200.00",
        "Onboarding session       2       450.00",
        "",
        "Total due                       1650.00",
    ]
    y = 300
    for line in body:
        draw.text((120, y), line, fill=(30, 30, 30), font=_font(34))
        y += 56

    # JPEG rather than PNG: PyMuPDF embeds the raster as-is, and a lossless page
    # scan would make the fixture several megabytes for no benefit.
    raster_path = path.with_suffix(".page.jpg")
    image.save(raster_path, format="JPEG", quality=72, optimize=True)

    document = fitz.open()
    page = document.new_page(width=595, height=842)
    page.insert_image(fitz.Rect(0, 0, 595, 842), filename=str(raster_path))
    document.save(path)
    document.close()
    raster_path.unlink(missing_ok=True)
    return path


def make_docx(path: Path) -> Path:
    """A DOCX with two heading levels, a bullet list and a 3x3 table."""
    from docx import Document

    document = Document()
    document.add_heading("Release Checklist", level=1)
    document.add_paragraph("This document tracks the 1.4.2 release.")
    document.add_heading("Blocking items", level=2)
    for item in ("Fix report export crash", "Update deployment runbook"):
        document.add_paragraph(item, style="List Bullet")

    table = document.add_table(rows=3, cols=3)
    headers = ("Item", "Owner", "Status")
    rows = (("Export crash", "Mai", "In progress"), ("Runbook", "Tom", "Done"))
    for column, text in enumerate(headers):
        table.cell(0, column).text = text
    for row_index, row in enumerate(rows, start=1):
        for column, text in enumerate(row):
            table.cell(row_index, column).text = text

    document.save(path)
    return path


def make_pptx(path: Path, image_path: Path) -> Path:
    """A two-slide deck; the second slide embeds an image."""
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Q1 Reliability Review"
    slide.placeholders[1].text = "Uptime 99.4%\nTwo Sev-2 incidents\nMTTR down to 34 minutes"

    second = presentation.slides.add_slide(presentation.slide_layouts[5])
    second.shapes.title.text = "Build failure example"
    second.shapes.add_picture(str(image_path), Inches(1), Inches(2), width=Inches(6))

    presentation.save(path)
    return path


def make_xlsx(path: Path) -> Path:
    """A workbook with two named sheets holding small tables."""
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Incidents"
    sheet.append(["Date", "Severity", "Component", "Minutes"])
    sheet.append(["2026-01-14", "Sev-2", "report-export", 41])
    sheet.append(["2026-02-02", "Sev-2", "auth-gateway", 27])

    second = workbook.create_sheet("Budget")
    second.append(["Item", "Cost"])
    second.append(["Monitoring", 240])
    second.append(["On-call", 900])

    workbook.save(path)
    return path


def make_plain_files(directory: Path) -> dict[str, Path]:
    files = {}

    txt = directory / "notes.txt"
    txt.write_text(
        "Release notes 1.4.2\n"
        "===================\n"
        "- Fixed report export crash on empty datasets.\n"
        "- Raised the upload limit to 50 MB.\n",
        encoding="utf-8",
    )
    files["txt"] = txt

    md = directory / "readme.md"
    md.write_text(
        "# Service overview\n\n"
        "The report service renders PDF exports.\n\n"
        "## Limits\n\n"
        "| Setting | Value |\n| --- | --- |\n| Max upload | 50 MB |\n",
        encoding="utf-8",
    )
    files["md"] = md

    csv = directory / "metrics.csv"
    csv.write_text(
        "date,requests,errors\n2026-01-01,15234,12\n2026-01-02,16110,3\n",
        encoding="utf-8",
    )
    files["csv"] = csv

    html = directory / "page.html"
    html.write_text(
        "<html><head><title>Runbook</title></head><body>"
        "<h1>Runbook</h1><p>Restart the worker pool.</p>"
        "<ul><li>Drain queue</li><li>Restart</li></ul>"
        "</body></html>",
        encoding="utf-8",
    )
    files["html"] = html

    latin = directory / "legacy_cp1252.txt"
    latin.write_bytes("Café déjà vu — legacy encoding\n".encode("cp1252"))
    files["cp1252"] = latin

    return files


def generate_all(directory: Path | None = None) -> dict[str, Path]:
    """Build every fixture and return a name -> path map."""
    target = directory or GENERATED
    target.mkdir(parents=True, exist_ok=True)

    paths: dict[str, Path] = {}
    paths["terminal_png"] = make_terminal_screenshot(target / "terminal_error.png")
    paths["ui_png"] = make_ui_screenshot(target / "ui_settings.png")
    paths["text_pdf"] = make_text_pdf(target / "deployment.pdf")
    paths["scanned_pdf"] = make_scanned_pdf(target / "invoice_scanned.pdf")
    paths["docx"] = make_docx(target / "checklist.docx")
    paths["pptx"] = make_pptx(target / "review.pptx", paths["terminal_png"])
    paths["xlsx"] = make_xlsx(target / "incidents.xlsx")
    paths.update(make_plain_files(target))
    return paths


if __name__ == "__main__":  # pragma: no cover - developer convenience
    for name, path in sorted(generate_all().items()):
        print(f"{name:14} {path}  ({path.stat().st_size:,} bytes)")
